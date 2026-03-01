import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def generate_board_deck(company_description: str, themes: list[dict], roles: list[str] = None) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = "Strategic Improvement Recommendations"
    p.font.size = Pt(36)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
    tf = sub.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI-Powered Employee Interview Analysis"
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER

    detail = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1.5))
    tf = detail.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Company: {company_description}"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = f"Interviews conducted: {len(roles or [])}"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

    # ── Slide 2: Top 3 Priorities (titles only) ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    p = title.text_frame.paragraphs[0]
    p.text = "Top 3 Strategic Priorities"
    p.font.size = Pt(28)
    p.font.bold = True

    for i, theme in enumerate(themes[:3]):
        y = Inches(1.5 + i * 1.8)
        box = slide.shapes.add_textbox(Inches(0.5), y, Inches(12), Inches(1.2))
        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"#{i + 1}  {theme['theme_name']}  ({theme['action_count']} mentions)"
        p.font.size = Pt(24)
        p.font.bold = True

    # ── Slide 3: Deep Dive on #1 Priority ──
    if themes:
        top = themes[0]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        p = title.text_frame.paragraphs[0]
        p.text = f"Deep Dive: {top['theme_name']}"
        p.font.size = Pt(28)
        p.font.bold = True

        # Explanation
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(1.5))
        tf = body.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = top['summary']
        p.font.size = Pt(14)

        # Up to 5 employee quotes
        quotes_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12), Inches(4))
        tf = quotes_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "What employees said:"
        p.font.size = Pt(16)
        p.font.bold = True

        for q in top.get('key_quotes', [])[:5]:
            p = tf.add_paragraph()
            p.text = f'"{q}"'
            p.font.size = Pt(12)
            p.font.italic = True
            p.space_before = Pt(6)

    # ── Slide 4: Methodology & Employee Roles ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    p = title.text_frame.paragraphs[0]
    p.text = "Methodology & Participants"
    p.font.size = Pt(28)
    p.font.bold = True

    # Methodology
    meth = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.5), Inches(5))
    tf = meth.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Approach"
    p.font.size = Pt(18)
    p.font.bold = True

    steps = [
        "1. AI-simulated confidential interviews with diverse employees",
        "2. Actionable recommendations extracted from each conversation",
        "3. Recommendations embedded and clustered via K-Means",
        "4. Themes ranked by frequency and summarized for the Board",
    ]
    for step in steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(12)
        p.space_before = Pt(4)

    # Unique interviewed professions
    unique_roles = list(dict.fromkeys(roles or []))

    role_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6), Inches(5))
    tf = role_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Interviewed Professions"
    p.font.size = Pt(18)
    p.font.bold = True

    if unique_roles:
        # Show up to 10 unique professions
        for role in unique_roles[:10]:
            p = tf.add_paragraph()
            p.text = f"• {role}"
            p.font.size = Pt(10)
            p.space_before = Pt(2)
        if len(unique_roles) > 10:
            p = tf.add_paragraph()
            p.text = f"… and {len(unique_roles) - 10} more unique professions"
            p.font.size = Pt(10)
            p.font.italic = True

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
