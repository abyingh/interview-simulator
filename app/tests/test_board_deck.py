import io
from django.test import TestCase
from pptx import Presentation

from app.board_deck import generate_board_deck


class GenerateBoardDeckTest(TestCase):

    def _load_prs(self, buf: io.BytesIO) -> Presentation:
        buf.seek(0)
        return Presentation(buf)

    def test_returns_bytes_io(self):
        themes = [
            {"theme_name": "Culture", "summary": "Improve culture", "key_quotes": ["q1", "q2"],
             "action_count": 10, "sample_actions": ["a1"]},
        ]
        result = generate_board_deck("Test Company", themes, roles=["CTO"])
        self.assertIsInstance(result, io.BytesIO)
        self.assertGreater(result.getvalue().__len__(), 0)

    def test_slide_count_with_themes(self):
        themes = [
            {"theme_name": "Theme 1", "summary": "Summary 1", "key_quotes": ["q1"], "action_count": 5},
            {"theme_name": "Theme 2", "summary": "Summary 2", "key_quotes": ["q2"], "action_count": 3},
            {"theme_name": "Theme 3", "summary": "Summary 3", "key_quotes": ["q3"], "action_count": 2},
        ]
        buf = generate_board_deck("Company X", themes, roles=["CTO", "VP Eng"])
        prs = self._load_prs(buf)
        # 4 slides: title, top 3 priorities, deep dive #1, methodology
        self.assertEqual(len(prs.slides), 4)

    def test_slide_count_no_themes(self):
        buf = generate_board_deck("Company X", [], roles=["CTO"])
        prs = self._load_prs(buf)
        # Without themes: title, top 3 priorities (empty), methodology = 3 slides (no deep dive)
        self.assertEqual(len(prs.slides), 3)

    def test_title_slide_content(self):
        themes = [
            {"theme_name": "T1", "summary": "S1", "key_quotes": ["q1"], "action_count": 5},
        ]
        buf = generate_board_deck("Acme Corp", themes, roles=["Engineer"])
        prs = self._load_prs(buf)

        title_slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in title_slide.shapes if shape.has_text_frame]
        full_text = ' '.join(texts)

        self.assertIn("Strategic Improvement Recommendations", full_text)
        self.assertIn("Acme Corp", full_text)

    def test_deep_dive_slide_content(self):
        themes = [
            {"theme_name": "Innovation Gap", "summary": "Need more R&D investment",
             "key_quotes": ["We lack innovation", "R&D budget is low"], "action_count": 8},
        ]
        buf = generate_board_deck("TechCo", themes, roles=["CTO"])
        prs = self._load_prs(buf)

        # Deep dive slide is the 3rd slide (index 2)
        deep_dive = prs.slides[2]
        texts = [shape.text_frame.text for shape in deep_dive.shapes if shape.has_text_frame]
        full_text = ' '.join(texts)

        self.assertIn("Innovation Gap", full_text)
        self.assertIn("Need more R&D investment", full_text)
        self.assertIn("We lack innovation", full_text)

    def test_methodology_slide_roles(self):
        roles = ["CTO", "VP Engineering", "Data Scientist", "CTO"]  # CTO duplicated
        themes = [
            {"theme_name": "T1", "summary": "S", "key_quotes": ["q"], "action_count": 5},
        ]
        buf = generate_board_deck("Company", themes, roles=roles)
        prs = self._load_prs(buf)

        methodology_slide = prs.slides[-1]
        texts = [shape.text_frame.text for shape in methodology_slide.shapes if shape.has_text_frame]
        full_text = ' '.join(texts)

        self.assertIn("CTO", full_text)
        self.assertIn("VP Engineering", full_text)
        self.assertIn("Data Scientist", full_text)

    def test_no_roles(self):
        themes = [
            {"theme_name": "T1", "summary": "S", "key_quotes": ["q"], "action_count": 5},
        ]
        buf = generate_board_deck("Company", themes, roles=None)
        self.assertIsInstance(buf, io.BytesIO)
        prs = self._load_prs(buf)
        self.assertEqual(len(prs.slides), 4)

    def test_many_key_quotes_truncated(self):
        themes = [
            {
                "theme_name": "Theme",
                "summary": "Summary",
                "key_quotes": [f"Quote {i}" for i in range(10)],  # 10 quotes, max shown is 5
                "action_count": 20,
            },
        ]
        buf = generate_board_deck("Company", themes, roles=["Role"])
        prs = self._load_prs(buf)

        deep_dive = prs.slides[2]
        texts = [shape.text_frame.text for shape in deep_dive.shapes if shape.has_text_frame]
        full_text = ' '.join(texts)

        # Quotes 0-4 should be present, quote 5+ should not
        for i in range(5):
            self.assertIn(f"Quote {i}", full_text)
        self.assertNotIn("Quote 5", full_text)
